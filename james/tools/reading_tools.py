"""Reading tools: extract text from PDFs, office documents, audio, and images.

These are the Level-1 capabilities GAIA-style benchmarks expect: most tasks
hand the agent a file (PDF, spreadsheet, audio clip, image, archive) and ask
a question about its contents. ``read_file`` only handles plain text, so this
module closes the multimodal gap.
"""

from __future__ import annotations

import base64
import csv
import io
from pathlib import Path

from ..config import settings
from ..core.workspace import resolve_workspace_path
from .base import ToolResult, tool

_MAX_FILE_BYTES = 25 * 1024 * 1024

_MIME = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".bmp": "image/bmp",
}

_TEXT_SUFFIXES = {
    ".txt",
    ".md",
    ".log",
    ".json",
    ".py",
    ".js",
    ".ts",
    ".html",
    ".htm",
    ".xml",
    ".yaml",
    ".yml",
    ".ini",
    ".cfg",
    ".toml",
    ".csv",
    ".tsv",
}


def _open_document(path: Path, max_chars: int) -> ToolResult:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return read_pdf.run(path=str(path), max_chars=max_chars)
    if suffix in (".docx",):
        try:
            import docx  # python-docx
        except ImportError:
            return ToolResult(
                ok=False,
                output="Reading .docx files requires the [docs] extra. "
                "Run: pip install 'james-assistant[docs]'",
            )
        document = docx.Document(str(path))
        parts = [p.text for p in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        return ToolResult(ok=True, output=("\n".join(parts)).strip()[:max_chars] or "(no text)")
    if suffix in (".pptx",):
        try:
            from pptx import Presentation
        except ImportError:
            return ToolResult(
                ok=False,
                output="Reading .pptx files requires the [docs] extra. "
                "Run: pip install 'james-assistant[docs]'",
            )
        presentation = Presentation(str(path))
        lines: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, 1):
            lines.append(f"--- Slide {slide_index} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs).strip()
                        if text:
                            lines.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        lines.append(" | ".join(cell.text.strip() for cell in row.cells))
        return ToolResult(ok=True, output="\n".join(lines)[:max_chars] or "(no text)")
    if suffix in (".xlsx", ".xlsm"):
        try:
            import openpyxl
        except ImportError:
            return ToolResult(
                ok=False,
                output="Reading .xlsx files requires the [docs] extra. "
                "Run: pip install 'james-assistant[docs]'",
            )
        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        rows: list[str] = []
        for sheet in workbook.worksheets:
            rows.append(f"--- Sheet: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                cells = ["" if v is None else str(v) for v in row]
                if any(cells):
                    rows.append(" | ".join(cells))
        workbook.close()
        return ToolResult(ok=True, output="\n".join(rows)[:max_chars] or "(empty workbook)")
    if suffix in (".csv", ".tsv"):
        delimiter = "\t" if suffix == ".tsv" else ","
        try:
            text = path.read_text(encoding="utf-8", errors="replace")
            rows = [list(r) for r in csv.reader(io.StringIO(text), delimiter=delimiter)]
            lines = [" | ".join(row) for row in rows]
            return ToolResult(ok=True, output="\n".join(lines)[:max_chars] or "(empty file)")
        except Exception as exc:
            return ToolResult(ok=False, output=f"Cannot read CSV: {exc}")
    if suffix in (".ods",):
        try:
            from odf.opendocument import load as odf_load
            from odf.table import Table, TableCell, TableRow
            from odf.teletype import extractText
        except ImportError:
            return ToolResult(
                ok=False,
                output="Reading .ods files requires the [docs] extra (odfpy). "
                "Run: pip install 'james-assistant[docs]'",
            )
        try:
            doc = odf_load(str(path))
            lines: list[str] = []
            for table in doc.spreadsheet.getElementsByType(Table):
                lines.append(f"--- Sheet: {table.getAttribute('name')} ---")
                for row in table.getElementsByType(TableRow):
                    cells: list[str] = []
                    for cell in row.getElementsByType(TableCell):
                        cells.append(extractText(cell).strip())
                    if any(cells):
                        lines.append(" | ".join(cells))
            return ToolResult(ok=True, output="\n".join(lines)[:max_chars] or "(empty spreadsheet)")
        except Exception as exc:
            return ToolResult(ok=False, output=f"Cannot read .ods file: {exc}")
    if suffix in _TEXT_SUFFIXES:
        try:
            return ToolResult(
                ok=True, output=path.read_text(encoding="utf-8", errors="replace")[:max_chars]
            )
        except Exception as exc:
            return ToolResult(ok=False, output=f"Cannot read file: {exc}")
    return ToolResult(
        ok=False,
        output=(
            f"Unsupported format: {suffix or '(none)'}. Supported: pdf, docx, pptx, "
            "xlsx, csv, tsv, ods, and plain text files."
        ),
    )


@tool(
    "read_pdf",
    "Extract the full text of a PDF file (pages joined with blank lines).",
    {
        "path": {"type": "string", "description": "Path to the PDF file."},
        "max_chars": {
            "type": "integer",
            "description": "Cap on returned characters (default 20000).",
        },
    },
    required=["path"],
)
def read_pdf(path: str, max_chars: int = 20000) -> ToolResult:
    p = resolve_workspace_path(path)
    if not p.exists():
        return ToolResult(ok=False, output=f"File not found: {path}")
    if p.stat().st_size > _MAX_FILE_BYTES:
        return ToolResult(ok=False, output=f"File too large: {p.stat().st_size} bytes")
    try:
        from pypdf import PdfReader
    except ImportError:
        return ToolResult(
            ok=False,
            output="PDF reading requires the [docs] extra (pypdf). "
            "Run: pip install 'james-assistant[docs]'",
        )
    try:
        reader = PdfReader(str(p))
        pages = [page.extract_text() or "" for page in reader.pages]
        text = "\n\n".join(pages).strip()
        return ToolResult(ok=True, output=text[:max_chars] or "(no extractable text)")
    except Exception as exc:
        return ToolResult(ok=False, output=f"PDF read failed: {exc}")


@tool(
    "read_document",
    "Read an office document (docx, pptx, xlsx, csv, tsv, ods) or plain text file and return its contents as text.",
    {
        "path": {"type": "string", "description": "Path to the document."},
        "max_chars": {
            "type": "integer",
            "description": "Cap on returned characters (default 20000).",
        },
    },
    required=["path"],
)
def read_document(path: str, max_chars: int = 20000) -> ToolResult:
    p = resolve_workspace_path(path)
    if not p.exists():
        return ToolResult(ok=False, output=f"File not found: {path}")
    return _open_document(p, max_chars)


@tool(
    "extract_audio_text",
    "Transcribe a local audio file (wav, mp3, m4a, ogg, flac) into text using Whisper.",
    {
        "path": {"type": "string", "description": "Path to the audio file."},
        "language": {
            "type": "string",
            "description": "Optional language code (e.g. 'en'); auto-detected when omitted.",
        },
    },
    required=["path"],
)
def extract_audio_text(path: str, language: str = "") -> ToolResult:
    p = resolve_workspace_path(path)
    if not p.exists():
        return ToolResult(ok=False, output=f"File not found: {path}")
    try:
        import whisper
    except ImportError:
        return ToolResult(
            ok=False,
            output="Audio transcription requires the [voice] extra (openai-whisper). "
            "Run: pip install 'james-assistant[voice]'",
        )
    try:
        model = whisper.load_model("base")
        kwargs = {"language": language} if language else {}
        result = model.transcribe(str(p), **kwargs)
        return ToolResult(
            ok=True, output=(result.get("text") or "").strip() or "(no speech detected)"
        )
    except Exception as exc:
        return ToolResult(ok=False, output=f"Transcription failed: {exc}")


@tool(
    "describe_image",
    "Analyze a local image with a vision-capable model and return a text description. Handles PNG, JPEG, WebP, GIF, BMP.",
    {
        "path": {"type": "string", "description": "Path to the image file."},
        "question": {
            "type": "string",
            "description": "Optional question about the image; otherwise a full description is returned.",
        },
    },
    required=["path"],
)
def describe_image(path: str, question: str = "") -> ToolResult:
    p = resolve_workspace_path(path)
    if not p.exists():
        return ToolResult(ok=False, output=f"File not found: {path}")
    if p.stat().st_size > 10_000_000:
        return ToolResult(ok=False, output=f"File too large: {p.stat().st_size} bytes (max 10MB)")
    mime = _MIME.get(p.suffix.lower(), "image/png")
    uri = f"data:{mime};base64,{base64.b64encode(p.read_bytes()).decode('ascii')}"
    prompt = question or (
        "Describe the contents of this image in detail. Include any text, numbers, "
        "names, objects, or relationships visible. Be precise and complete."
    )
    try:
        from ..llm.factory import build_provider

        provider = build_provider(settings.llm)
        model = settings.assistant.vision_model or None
        response = provider.chat(
            [{"role": "user", "content": prompt}],
            images=[uri],
            model=model,
        )
        return ToolResult(ok=True, output=(response.content or "").strip() or "(no response)")
    except Exception as exc:
        return ToolResult(
            ok=False,
            output=f"Image analysis failed: {exc} (the configured model may not support vision).",
        )
