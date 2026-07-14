"""Document generation tools: Word (.docx), PowerPoint (.pptx) and PDF."""
from __future__ import annotations

import json
from pathlib import Path

from ..config import settings
from .base import Tool, ToolResult, tool


def _out(name: str) -> Path:
    p = settings.assistant.workspace_dir / name
    p.parent.mkdir(parents=True, exist_ok=True)
    return p


@tool(
    "create_word_document",
    "Create a formatted Microsoft Word document (.docx) with a title and structured sections. Returns the saved path.",
    {
        "filename": {"type": "string", "description": "Output file name, e.g. report.docx"},
        "title": {"type": "string", "description": "Document title."},
        "sections": {
            "type": "array",
            "description": "List of sections, each with 'heading' and 'body'.",
            "items": {
                "type": "object",
                "properties": {
                    "heading": {"type": "string"},
                    "body": {"type": "string"},
                },
            },
        },
    },
    required=["filename", "title", "sections"],
)
def create_word_document(filename: str, title: str, sections: list) -> ToolResult:
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=0)
    for sec in sections or []:
        doc.add_heading(sec.get("heading", ""), level=1)
        doc.add_paragraph(sec.get("body", ""))
    p = _out(filename)
    doc.save(str(p))
    return ToolResult(ok=True, output=f"Created Word document at {p}")


@tool(
    "create_powerpoint",
    "Create a Microsoft PowerPoint presentation (.pptx) with title slides and bullet points.",
    {
        "filename": {"type": "string", "description": "Output file name, e.g. deck.pptx"},
        "title": {"type": "string", "description": "Presentation title (first slide)."},
        "slides": {
            "type": "array",
            "description": "Slides, each with 'title' and 'bullets' (list of strings).",
            "items": {
                "type": "object",
                "properties": {
                    "title": {"type": "string"},
                    "bullets": {"type": "array", "items": {"type": "string"}},
                },
            },
        },
    },
    required=["filename", "title", "slides"],
)
def create_powerpoint(filename: str, title: str, slides: list) -> ToolResult:
    from pptx import Presentation

    prs = Presentation()
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = title

    for s in slides or []:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s.get("title", "")
        body = slide.placeholders[1].text_frame
        for i, b in enumerate(s.get("bullets", [])):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = f"• {b}"
    p = _out(filename)
    prs.save(str(p))
    return ToolResult(ok=True, output=f"Created PowerPoint at {p}")


@tool(
    "create_pdf",
    "Create a PDF document from a title and sections of text.",
    {
        "filename": {"type": "string", "description": "Output file name, e.g. report.pdf"},
        "title": {"type": "string", "description": "Document title."},
        "sections": {
            "type": "array",
            "description": "List of sections with 'heading' and 'body'.",
            "items": {
                "type": "object",
                "properties": {
                    "heading": {"type": "string"},
                    "body": {"type": "string"},
                },
            },
        },
    },
    required=["filename", "title", "sections"],
)
def create_pdf(filename: str, title: str, sections: list) -> ToolResult:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    p = _out(filename)
    doc = SimpleDocTemplate(str(p), pagesize=A4)
    styles = getSampleStyleSheet()
    story = [Paragraph(title, styles["Title"]), Spacer(1, 12)]
    for sec in sections or []:
        story.append(Paragraph(sec.get("heading", ""), styles["Heading2"]))
        story.append(Paragraph(sec.get("body", ""), styles["BodyText"]))
        story.append(Spacer(1, 8))
    doc.build(story)
    return ToolResult(ok=True, output=f"Created PDF at {p}")
