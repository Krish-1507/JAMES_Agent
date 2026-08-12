"""Windows app automation — drive the user's real Office applications via COM.

Outlook (inbox/email/calendar), Excel, Word and PowerPoint are controlled
through ``pywin32`` COM on Windows when Microsoft Office is installed. The
file-based tools (``excel_read_cells``, ``excel_write_cells``,
``word_read_document``) transparently fall back to openpyxl / python-docx so
they also work on other platforms.

Every tool fails with a clear message when the platform or dependency is
missing — none of them ever blocks the agent loop.
"""

from __future__ import annotations

import json
import sys
from pathlib import Path

from ..core.workspace import resolve_workspace_path
from .base import ToolResult, tool


def _com_dispatch(progid: str):
    """Return a live COM object or None (with the reason) when unavailable."""
    if sys.platform != "win32":
        raise RuntimeError("Office COM automation requires Windows.")
    try:
        import win32com.client  # type: ignore[import-not-found]

        return win32com.client.Dispatch(progid)
    except Exception as exc:
        raise RuntimeError(f"Office COM dispatch failed ({progid}): {exc}") from exc


def _guard_office(feature: str) -> None:
    """Raise if the platform/dependency for COM Office automation is missing."""
    if sys.platform != "win32":
        raise RuntimeError(
            f"{feature} uses the installed Microsoft Office via COM, which is "
            "Windows-only. On other platforms use the file-based document tools."
        )
    try:
        import win32com.client  # noqa: F401
    except ImportError as exc:
        raise RuntimeError(f"{feature} needs pywin32: pip install pywin32") from exc


def _as_json(value: str, name: str) -> list:
    try:
        parsed = json.loads(value or "[]")
    except json.JSONDecodeError as exc:
        raise ValueError(f"{name} must be valid JSON, got: {value[:80]}") from exc
    if not isinstance(parsed, list):
        raise ValueError(f"{name} must be a JSON list")
    return parsed


@tool(
    "outlook_read_inbox",
    "List recent emails from the user's Outlook inbox (Windows + Microsoft Office).",
    {
        "count": {"type": "integer", "description": "How many emails to return (default 10)."},
        "unread_only": {
            "type": "boolean",
            "description": "Only show unread messages (default false).",
        },
    },
)
def outlook_read_inbox(count: int = 10, unread_only: bool = False) -> ToolResult:
    try:
        _guard_office("outlook_read_inbox")
        app = _com_dispatch("Outlook.Application")
        namespace = app.GetNamespace("MAPI")
        inbox = namespace.GetDefaultFolder(6)  # olFolderInbox
        items = []
        for msg in list(inbox.Items)[: max(1, int(count))]:
            try:
                if unread_only and not msg.UnRead:
                    continue
                items.append(
                    {
                        "from": str(getattr(msg, "SenderName", "") or ""),
                        "subject": str(getattr(msg, "Subject", "") or ""),
                        "received": str(getattr(msg, "ReceivedTime", "") or ""),
                        "unread": bool(getattr(msg, "UnRead", False)),
                    }
                )
            except Exception:
                continue
        if not items:
            return ToolResult(ok=True, output="No emails found.")
        return ToolResult(
            ok=True,
            output=json.dumps(items, indent=2, ensure_ascii=False),
            data=items,
        )
    except Exception as exc:
        return ToolResult(ok=False, output=str(exc))


@tool(
    "outlook_send_email",
    "Send an email through the user's Outlook (Windows + Microsoft Office).",
    {
        "to": {
            "type": "string",
            "description": "Recipient email address (comma-separated for several).",
        },
        "subject": {"type": "string", "description": "Email subject."},
        "body": {"type": "string", "description": "Plain-text email body."},
        "attachments": {
            "type": "array",
            "items": {"type": "string"},
            "description": "Optional workspace file paths to attach.",
        },
    },
    required=["to", "subject"],
)
def outlook_send_email(
    to: str, subject: str, body: str = "", attachments: list | None = None
) -> ToolResult:
    try:
        _guard_office("outlook_send_email")
        app = _com_dispatch("Outlook.Application")
        mail = app.CreateItem(0)  # olMailItem
        mail.To = to
        mail.Subject = subject
        mail.Body = body
        for attachment in attachments or []:
            path = resolve_workspace_path(attachment, allow_root=False)
            if not path.exists():
                return ToolResult(ok=False, output=f"Attachment not found: {path}")
            mail.Attachments.Add(str(path))
        mail.Send()
        return ToolResult(ok=True, output=f"Email sent to {to}.")
    except Exception as exc:
        return ToolResult(ok=False, output=str(exc))


@tool(
    "outlook_create_event",
    "Create a calendar event in the user's Outlook (Windows + Microsoft Office).",
    {
        "subject": {"type": "string", "description": "Event title."},
        "start": {"type": "string", "description": "Start time, e.g. '2026-08-12T09:00:00'."},
        "end": {"type": "string", "description": "End time, e.g. '2026-08-12T10:00:00'."},
        "location": {"type": "string", "description": "Optional location."},
        "body": {"type": "string", "description": "Optional event details."},
    },
    required=["subject", "start", "end"],
)
def outlook_create_event(
    subject: str, start: str, end: str, location: str = "", body: str = ""
) -> ToolResult:
    try:
        _guard_office("outlook_create_event")
        app = _com_dispatch("Outlook.Application")
        event = app.CreateItem(1)  # olAppointmentItem
        event.Subject = subject
        event.Start = start
        event.End = end
        if location:
            event.Location = location
        if body:
            event.Body = body
        event.Save()
        return ToolResult(ok=True, output=f"Calendar event '{subject}' saved.")
    except Exception as exc:
        return ToolResult(ok=False, output=str(exc))


def _excel_com_read(path: Path, sheet: str, cell_range: str) -> list[list[str]] | None:
    app = _com_dispatch("Excel.Application")
    workbook = app.Workbooks.Open(str(path))
    try:
        ws = workbook.Worksheets(sheet)
        values = ws.Range(cell_range).Value
        if values is None:
            return []
        if not isinstance(values, tuple):  # single cell
            return [[str(values)]]
        return [[str(cell) if cell is not None else "" for cell in row] for row in values]
    finally:
        workbook.Close(False)
        app.Quit()


def _excel_openpyxl_read(path: Path, sheet: str, cell_range: str) -> list[list[str]]:
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), data_only=True)
    ws = workbook[sheet] if sheet in workbook.sheetnames else workbook.active
    rows = []
    for row in ws[cell_range]:
        rows.append([str(cell.value) if cell.value is not None else "" for cell in row])
    return rows


@tool(
    "excel_read_cells",
    "Read a cell range from an Excel workbook. Uses live Excel (Windows) or openpyxl.",
    {
        "path": {"type": "string", "description": "Workbook path (workspace-relative)."},
        "sheet": {"type": "string", "description": "Sheet name (default first sheet)."},
        "range": {"type": "string", "description": "Cell range, e.g. 'A1:C10'."},
    },
    required=["path"],
)
def excel_read_cells(path: str, sheet: str = "", range: str = "A1:Z100") -> ToolResult:
    try:
        p = resolve_workspace_path(path, allow_root=False)
        if not p.exists():
            return ToolResult(ok=False, output=f"File not found: {path}")
        sheet_name = sheet or "Sheet1"
        rows = []
        try:
            rows = _excel_com_read(p, sheet_name, range)
        except Exception:
            rows = _excel_openpyxl_read(p, sheet_name, range)
        if not rows:
            return ToolResult(ok=True, output="The range is empty.")
        return ToolResult(
            ok=True,
            output=json.dumps(rows, ensure_ascii=False),
            data=rows,
        )
    except Exception as exc:
        return ToolResult(ok=False, output=str(exc))


def _excel_com_write(path: Path, rows: list[list], sheet: str, start_cell: str) -> None:
    app = _com_dispatch("Excel.Application")
    try:
        workbook = app.Workbooks.Open(str(path)) if path.exists() else app.Workbooks.Add()
        try:
            ws = (
                workbook.Worksheets(sheet)
                if sheet in [s.Name for s in workbook.Worksheets]
                else workbook.Worksheets.Add()
            )
            ws.Name = sheet
        except Exception:
            ws = workbook.Worksheets(1)
        top = ws.Range(start_cell)
        for row_offset, row in enumerate(rows):
            for col_offset, value in enumerate(row):
                top.Offset(row_offset, col_offset).Value = value
        if not path.exists():
            workbook.SaveAs(str(path))
        else:
            workbook.Save()
    finally:
        workbook.Close(False)
        app.Quit()


def _excel_openpyxl_write(path: Path, rows: list[list], sheet: str, start_cell: str) -> None:
    from openpyxl import Workbook, load_workbook

    workbook = load_workbook(str(path)) if path.exists() else Workbook()
    ws = workbook[sheet] if sheet in workbook.sheetnames else workbook.create_sheet(sheet)
    for row_offset, row in enumerate(rows):
        for col_offset, value in enumerate(row):
            cell = ws[start_cell]
            ws.cell(row=cell.row + row_offset, column=cell.column + col_offset, value=value)
    workbook.save(str(path))


@tool(
    "excel_write_cells",
    "Write values into an Excel workbook, creating the file if needed. "
    "Uses live Excel (Windows) or openpyxl.",
    {
        "path": {"type": "string", "description": "Workbook path (workspace-relative)."},
        "data": {
            "type": "string",
            "description": 'Rows as a JSON list of lists, e.g. [["Name","Score"],["Ada",9]].',
        },
        "sheet": {"type": "string", "description": "Sheet name (default 'Sheet1')."},
        "start_cell": {"type": "string", "description": "Top-left cell, e.g. 'A1'."},
    },
    required=["path", "data"],
)
def excel_write_cells(
    path: str, data: str, sheet: str = "Sheet1", start_cell: str = "A1"
) -> ToolResult:
    try:
        p = resolve_workspace_path(path, allow_root=False)
        rows = _as_json(data, "data")
        if any(not isinstance(r, list) for r in rows):
            return ToolResult(ok=False, output="data must be a list of rows (lists of values).")
        try:
            _excel_com_write(p, rows, sheet, start_cell)
        except Exception:
            _excel_openpyxl_write(p, rows, sheet, start_cell)
        return ToolResult(ok=True, output=f"Wrote {len(rows)} row(s) to {p}.")
    except Exception as exc:
        return ToolResult(ok=False, output=str(exc))


@tool(
    "word_read_document",
    "Extract the text content of a Word document. Uses live Word (Windows) or python-docx.",
    {"path": {"type": "string", "description": "Document path (workspace-relative)."}},
    required=["path"],
)
def word_read_document(path: str) -> ToolResult:
    try:
        p = resolve_workspace_path(path, allow_root=False)
        if not p.exists():
            return ToolResult(ok=False, output=f"File not found: {path}")
        try:
            _guard_office("word_read_document")
            app = _com_dispatch("Word.Application")
            app.Visible = False
            try:
                doc = app.Documents.Open(str(p))
                text = doc.Content.Text
                doc.Close(False)
                return ToolResult(ok=True, output=text.strip()[:100_000])
            finally:
                app.Quit()
        except Exception:
            from docx import Document

            doc = Document(str(p))
            return ToolResult(
                ok=True,
                output="\n".join(par.text for par in doc.paragraphs).strip()[:100_000],
            )
    except Exception as exc:
        return ToolResult(ok=False, output=str(exc))


@tool(
    "powerpoint_create",
    "Create a PowerPoint presentation from a JSON outline. Uses live PowerPoint (Windows).",
    {
        "title": {"type": "string", "description": "Presentation title (slide 1)."},
        "slides": {
            "type": "string",
            "description": 'JSON list of slides: [{"title": str, "bullets": [str, ...]}, ...].',
        },
        "save_path": {
            "type": "string",
            "description": "Destination .pptx path (workspace-relative, default slides.pptx).",
        },
    },
    required=["title"],
)
def powerpoint_create(title: str, slides: str = "[]", save_path: str = "slides.pptx") -> ToolResult:
    try:
        outline = _as_json(slides, "slides")
        p = resolve_workspace_path(save_path, allow_root=False)
        try:
            _guard_office("powerpoint_create")
        except Exception as exc:
            return _powerpoint_pptx_fallback(title, outline, p, reason=str(exc))
        app = _com_dispatch("PowerPoint.Application")
        try:
            presentation = app.Presentations.Add()
            slide = presentation.Slides.Add(1, 1)  # ppLayoutTitle
            slide.Shapes.Title.TextFrame.TextRange.Text = title
            for index, item in enumerate(outline, start=2):
                slide_title = str(item.get("title", f"Slide {index - 1}"))
                bullets = [str(b) for b in (item.get("bullets") or [])]
                new_slide = presentation.Slides.Add(index, 1)
                new_slide.Shapes.Title.TextFrame.TextRange.Text = slide_title
                body = new_slide.Shapes.Placeholders(2)
                text_frame = body.TextFrame
                text_frame.TextRange.Text = "\n".join(bullets) if bullets else ""
            presentation.SaveAs(str(p))
            presentation.Close()
            return ToolResult(ok=True, output=f"Presentation saved to {p}.")
        finally:
            app.Quit()
    except Exception as exc:
        return ToolResult(ok=False, output=str(exc))


def _powerpoint_pptx_fallback(title: str, outline: list, p: Path, reason: str) -> ToolResult:
    """Build the deck with python-pptx when COM/Office is unavailable."""
    try:
        from pptx import Presentation
    except ImportError:
        return ToolResult(
            ok=False,
            output=(
                f"{reason} Install the [docs] extra (python-pptx) for a "
                "portable .pptx fallback that works on every platform."
            ),
        )
    prs = Presentation()
    cover = prs.slides.add_slide(prs.slide_layouts[0])
    cover.shapes.title.text = title
    for item in outline:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(item.get("title", ""))
        body = slide.placeholders[1].text_frame
        for i, b in enumerate([str(b) for b in (item.get("bullets") or [])]):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = f"• {b}"
    prs.save(str(p))
    return ToolResult(ok=True, output=f"Presentation saved to {p} (python-pptx fallback).")
