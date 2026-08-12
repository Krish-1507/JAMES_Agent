"""Phase-5 tests for document generation tools and voice provider factories."""

from __future__ import annotations

from pathlib import Path

import pytest

from james.tools import document_tools as doc_tools

# ---------------------------------------------------------------------------
# document generation
# ---------------------------------------------------------------------------


def test_create_word_document(monkeypatch: pytest.MonkeyPatch, isolated_workspace: Path) -> None:
    result = doc_tools.create_word_document.run(
        filename="report.docx",
        title="Quarterly Report",
        sections=[{"heading": "Summary", "body": "It went well."}],
    )
    assert result.ok is True
    path = isolated_workspace / "report.docx"
    assert path.exists()
    from docx import Document

    doc = Document(str(path))
    assert doc.paragraphs[0].text == "Quarterly Report"


def test_create_powerpoint(monkeypatch: pytest.MonkeyPatch, isolated_workspace: Path) -> None:
    result = doc_tools.create_powerpoint.run(
        filename="deck.pptx", title="Demo", slides=[{"title": "Slide One", "bullets": ["a", "b"]}]
    )
    assert result.ok is True
    path = isolated_workspace / "deck.pptx"
    assert path.exists()
    from pptx import Presentation

    prs = Presentation(str(path))
    assert len(prs.slides) == 2
    assert prs.slides[0].shapes.title.text == "Demo"


def test_create_pdf(monkeypatch: pytest.MonkeyPatch, isolated_workspace: Path) -> None:
    result = doc_tools.create_pdf.run(
        filename="brief.pdf", title="Brief", sections=[{"heading": "Intro", "body": "Hello world."}]
    )
    assert result.ok is True
    path = isolated_workspace / "brief.pdf"
    assert path.exists()
    assert path.read_bytes().startswith(b"%PDF")


def test_create_word_outside_workspace_rejected(
    monkeypatch: pytest.MonkeyPatch, isolated_workspace: Path
) -> None:
    result = doc_tools.create_word_document.run(
        filename="../escape.docx", title="Title", sections=[{"heading": "h", "body": "b"}]
    )
    assert result.ok is False


def test_document_tools_registered() -> None:
    from james.tools.registry import ALL_TOOLS

    names = {t.name for t in ALL_TOOLS}
    assert {"create_word_document", "create_powerpoint", "create_pdf"} <= names


# ---------------------------------------------------------------------------
# STT / TTS provider factories
# ---------------------------------------------------------------------------


def test_build_stt_none_returns_text_stt(monkeypatch: pytest.MonkeyPatch) -> None:
    from james.config import settings
    from james.voice.stt import TextSTT, build_stt

    monkeypatch.setattr(settings.voice, "enabled", False)
    provider = build_stt(settings.voice)
    assert isinstance(provider, TextSTT)


def test_build_stt_whisper_api(monkeypatch: pytest.MonkeyPatch) -> None:
    from james.config import settings
    from james.voice.stt import WhisperApiSTT, build_stt

    monkeypatch.setattr(settings.voice, "enabled", True)
    monkeypatch.setattr(settings.voice, "stt_provider", "whisper_api")
    monkeypatch.setattr(settings.voice, "whisper_api_key", "sk-test")
    provider = build_stt(settings.voice)
    assert isinstance(provider, WhisperApiSTT)


def test_build_stt_import_error_falls_back_to_text(monkeypatch: pytest.MonkeyPatch) -> None:
    from james.config import settings
    from james.voice.stt import TextSTT, build_stt

    monkeypatch.setattr(settings.voice, "enabled", True)
    monkeypatch.setattr(settings.voice, "stt_provider", "whisper_local")

    def boom(*a, **k):
        raise ImportError("whisper not installed")

    monkeypatch.setattr("james.voice.stt.WhisperLocalSTT.__init__", boom)
    provider = build_stt(settings.voice)
    assert isinstance(provider, TextSTT)


def test_whisper_api_listen_roundtrip(monkeypatch: pytest.MonkeyPatch) -> None:
    from james.voice.stt import WhisperApiSTT

    class FakeAudio:
        def get_wav_data(self) -> bytes:
            return b"WAV"

    class FakeClient:
        def __init__(self):
            self.audio = self

        @property
        def transcriptions(self):
            return self

        def create(self, model, file):
            class Resp:
                text = "hello world"

            return Resp()

    fake_client = FakeClient()
    provider = WhisperApiSTT.__new__(WhisperApiSTT)
    provider.client = fake_client
    provider.mic_index = None
    monkeypatch.setattr("james.voice.stt._record", lambda mic_index=None: FakeAudio())
    assert provider.listen() == "hello world"


def test_whisper_api_listen_empty_audio(monkeypatch: pytest.MonkeyPatch) -> None:
    from james.voice.stt import WhisperApiSTT

    provider = WhisperApiSTT.__new__(WhisperApiSTT)
    provider.client = object()
    provider.mic_index = None
    monkeypatch.setattr("james.voice.stt._record", lambda mic_index=None: None)
    assert provider.listen() == ""


def test_text_stt_listen_typed_input(monkeypatch: pytest.MonkeyPatch) -> None:
    from james.voice.stt import TextSTT

    monkeypatch.setattr("builtins.input", lambda prompt: "hello")
    assert TextSTT().listen() == "hello"
    monkeypatch.setattr("builtins.input", lambda prompt: (_ for _ in ()).throw(EOFError()))
    assert TextSTT().listen() == ""


def test_build_tts_none_returns_none_tts(monkeypatch: pytest.MonkeyPatch) -> None:
    from james.config import settings
    from james.voice.tts import NoneTTS, build_tts

    monkeypatch.setattr(settings.voice, "enabled", False)
    provider = build_tts(settings.voice)
    assert isinstance(provider, NoneTTS)


def test_build_tts_edge_missing_dependency_falls_back(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from james.config import settings
    from james.voice.tts import NoneTTS, build_tts

    monkeypatch.setattr(settings.voice, "enabled", True)
    monkeypatch.setattr(settings.voice, "tts_provider", "edge")

    def boom(*a, **k):
        raise ImportError("edge_tts missing")

    monkeypatch.setattr("james.voice.tts.EdgeTTS.__init__", boom)
    monkeypatch.setattr("james.voice.tts.Pyttsx3TTS.__init__", boom)
    provider = build_tts(settings.voice)
    assert isinstance(provider, NoneTTS)


def test_none_tts_prints(monkeypatch: pytest.MonkeyPatch) -> None:
    from james.voice.tts import NoneTTS

    captured: list[str] = []
    monkeypatch.setattr("builtins.print", lambda text: captured.append(text))
    NoneTTS().speak("hi")
    assert captured == ["hi"]


def test_tts_audio_helpers(monkeypatch: pytest.MonkeyPatch) -> None:
    import os
    import tempfile

    from james.voice import tts

    with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as handle:
        path = handle.name
    try:
        tts._play_audio(path)  # no ffplay/open/xdg-open in CI; must not raise
        assert os.path.exists(path)
    finally:
        os.unlink(path)
